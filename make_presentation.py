"""
Generates TSP-ACO presentation (tsp_aco_presentation.pptx).
Run with:  python make_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x1E, 0x1E, 0x2E)   # dark navy
C_ACCENT   = RGBColor(0x89, 0xB4, 0xFA)   # soft blue
C_ACCENT2  = RGBColor(0xA6, 0xE3, 0xA1)   # soft green
C_WARN     = RGBColor(0xF3, 0x8B, 0xA8)   # soft red/pink
C_TEXT     = RGBColor(0xCD, 0xD6, 0xF4)   # light lavender
C_SUBTEXT  = RGBColor(0xA6, 0xAD, 0xC8)   # muted lavender
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_SURFACE  = RGBColor(0x31, 0x32, 0x44)   # card background
C_SURFACE2 = RGBColor(0x45, 0x47, 0x5A)   # lighter card

W  = Inches(13.33)   # widescreen width
H  = Inches(7.5)     # widescreen height

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor = C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def label(slide, text, x, y, w, h,
          size=Pt(18), bold=False, color=C_TEXT,
          align=PP_ALIGN.LEFT, wrap=True):
    tf_box = slide.shapes.add_textbox(x, y, w, h)
    tf = tf_box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    return tf_box


def heading(slide, text, y=Inches(0.35), size=Pt(32)):
    label(slide, text,
          x=Inches(0.55), y=y, w=Inches(12.2), h=Inches(0.7),
          size=size, bold=True, color=C_ACCENT)


def subheading(slide, text, y=Inches(1.1)):
    label(slide, text,
          x=Inches(0.55), y=y, w=Inches(12.2), h=Inches(0.45),
          size=Pt(18), color=C_SUBTEXT)


def hline(slide, y):
    line_box = box(slide, Inches(0.55), y, Inches(12.2), Pt(2),
                   fill_color=C_ACCENT)
    return line_box


def card(slide, x, y, w, h, color=C_SURFACE):
    return box(slide, x, y, w, h, fill_color=color)


def bullet_list(slide, items, x, y, w,
                size=Pt(17), color=C_TEXT, spacing=Inches(0.42)):
    for i, item in enumerate(items):
        label(slide, item,
              x=x, y=y + i * spacing, w=w, h=spacing,
              size=size, color=color)


def table_shape(slide, rows, cols, x, y, w, h,
                header_color=C_ACCENT, header_text_color=C_BG,
                row_colors=(C_SURFACE, C_SURFACE2),
                text_color=C_TEXT, font_size=Pt(14)):
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table
    col_w = w // cols
    for c in range(cols):
        tbl.columns[c].width = col_w

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            tf = cell.text_frame
            tf.word_wrap = True
            # background
            fill = cell.fill
            fill.solid()
            if r == 0:
                fill.fore_color.rgb = header_color
            else:
                fill.fore_color.rgb = row_colors[(r - 1) % 2]
            # text placeholder
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.font.size  = font_size
            run.font.bold  = (r == 0)
            run.font.color.rgb = header_text_color if r == 0 else text_color
    return tbl


def set_cell(tbl, r, c, text, bold=False, align=PP_ALIGN.CENTER, size=None):
    cell = tbl.cell(r, c)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.bold = bold
    if size:
        run.font.size = size


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank_slide(prs)
    bg(s)

    # Accent bar on left
    box(s, 0, 0, Inches(0.18), H, fill_color=C_ACCENT)

    # Large title
    label(s, "Traveling Salesman Problem",
          x=Inches(0.55), y=Inches(1.8), w=Inches(12.2), h=Inches(1.1),
          size=Pt(48), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    label(s, "Ant Colony Optimization  |  CPU vs CUDA Parallelization",
          x=Inches(0.55), y=Inches(3.0), w=Inches(12.2), h=Inches(0.6),
          size=Pt(24), color=C_ACCENT, align=PP_ALIGN.CENTER)

    label(s, "Performance benchmark on two machines",
          x=Inches(0.55), y=Inches(3.75), w=Inches(12.2), h=Inches(0.45),
          size=Pt(18), color=C_SUBTEXT, align=PP_ALIGN.CENTER)

    # Bottom bar
    box(s, 0, Inches(6.9), W, Inches(0.6), fill_color=C_SURFACE)
    label(s, "TrafficSystem  |  2025",
          x=Inches(0.55), y=Inches(6.95), w=Inches(12.2), h=Inches(0.4),
          size=Pt(14), color=C_SUBTEXT, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "The Problem: Traveling Salesman")
    hline(s, Inches(1.1))

    # Left column — description
    label(s, "Given a set of cities, find the shortest possible route that visits each city exactly once and returns to the starting point.",
          x=Inches(0.55), y=Inches(1.35), w=Inches(5.8), h=Inches(1.2),
          size=Pt(17), color=C_TEXT)

    bullets_l = [
        "  NP-hard — no known polynomial-time exact solution",
        "  50 cities  ->  49! / 2  ~  3 x 10^62 possible routes",
        "  Exact methods are impractical for large instances",
        "  Heuristics and metaheuristics are used in practice",
    ]
    bullet_list(s, bullets_l,
                x=Inches(0.55), y=Inches(2.65), w=Inches(5.8),
                size=Pt(16), spacing=Inches(0.46))

    # Right column — city instance info card
    card(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.6))
    label(s, "Our Instance",
          x=Inches(7.2), y=Inches(1.5), w=Inches(5.4), h=Inches(0.45),
          size=Pt(19), bold=True, color=C_ACCENT)

    instance_data = [
        ("Cities",           "50"),
        ("Coordinate range", "5 – 1,605  (x)"),
        ("",                 "5 – 1,175  (y)"),
        ("Distance metric",  "Euclidean"),
        ("Source",           "cities.txt"),
    ]
    for i, (k, v) in enumerate(instance_data):
        y = Inches(2.1) + i * Inches(0.52)
        label(s, k, x=Inches(7.3), y=y, w=Inches(2.2), h=Inches(0.45),
              size=Pt(15), color=C_SUBTEXT)
        label(s, v, x=Inches(9.6), y=y, w=Inches(2.8), h=Inches(0.45),
              size=Pt(15), bold=True, color=C_TEXT)


def slide_aco_intro(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "The Method: Ant Colony Optimization")
    hline(s, Inches(1.1))

    label(s, "A bio-inspired metaheuristic introduced by Marco Dorigo (1992), based on the foraging behavior of real ant colonies.",
          x=Inches(0.55), y=Inches(1.3), w=Inches(12.2), h=Inches(0.8),
          size=Pt(17), color=C_TEXT)

    # Three concept cards
    concepts = [
        (C_ACCENT,  "Stigmergy",
         "Ants communicate indirectly by depositing pheromone on the paths they traverse."),
        (C_ACCENT2, "Positive Feedback",
         "Shorter paths accumulate pheromone faster, attracting more ants and reinforcing the best routes."),
        (C_WARN,    "Evaporation",
         "Pheromone evaporates over time, preventing premature convergence to suboptimal solutions."),
    ]
    for i, (color, title, desc) in enumerate(concepts):
        cx = Inches(0.55) + i * Inches(4.3)
        card(s, cx, Inches(2.3), Inches(4.0), Inches(2.4))
        box(s, cx, Inches(2.3), Inches(4.0), Inches(0.08), fill_color=color)
        label(s, title,
              x=cx + Inches(0.15), y=Inches(2.45), w=Inches(3.7), h=Inches(0.45),
              size=Pt(18), bold=True, color=color)
        label(s, desc,
              x=cx + Inches(0.15), y=Inches(2.95), w=Inches(3.7), h=Inches(1.5),
              size=Pt(15), color=C_TEXT)

    label(s, "The algorithm converges toward near-optimal solutions without exhaustive enumeration.",
          x=Inches(0.55), y=Inches(5.1), w=Inches(12.2), h=Inches(0.5),
          size=Pt(17), color=C_SUBTEXT, align=PP_ALIGN.CENTER)


def slide_aco_algorithm(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Algorithm: How It Works")
    hline(s, Inches(1.1))

    steps = [
        ("1", "Initialize",
         "Set pheromone level t0 on every edge. Place ants at random starting cities."),
        ("2", "Build Tours",
         "Each ant builds a complete tour. At each step, the next city is chosen probabilistically:   "
         "P(i->j)  proportional to  t(i,j)^alpha  *  (1/d(i,j))^beta"),
        ("3", "Evaluate",
         "Compute the tour length for each ant. Track the global best solution found so far."),
        ("4", "Evaporate",
         "Reduce all pheromone levels:  t(i,j)  <-  t(i,j) * (1 - rho)"),
        ("5", "Deposit",
         "Each ant adds pheromone to its tour edges:  delta(t)  =  Q / tour_length"),
        ("6", "Repeat",
         "Go back to step 2 until the maximum number of iterations is reached."),
    ]

    for i, (num, title, desc) in enumerate(steps):
        col = i % 3
        row = i // 3
        cx = Inches(0.45) + col * Inches(4.25)
        cy = Inches(1.3)  + row * Inches(2.7)

        card(s, cx, cy, Inches(4.05), Inches(2.45))
        # Number badge
        box(s, cx, cy, Inches(0.5), Inches(0.5), fill_color=C_ACCENT)
        label(s, num, x=cx, y=cy, w=Inches(0.5), h=Inches(0.5),
              size=Pt(16), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        label(s, title,
              x=cx + Inches(0.55), y=cy + Inches(0.05), w=Inches(3.4), h=Inches(0.42),
              size=Pt(16), bold=True, color=C_ACCENT)
        label(s, desc,
              x=cx + Inches(0.1), y=cy + Inches(0.55), w=Inches(3.85), h=Inches(1.7),
              size=Pt(13), color=C_TEXT)


def slide_cuda(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "CUDA Parallelization")
    hline(s, Inches(1.1))

    # CPU side
    card(s, Inches(0.45), Inches(1.3), Inches(5.7), Inches(5.6), color=C_SURFACE)
    box(s, Inches(0.45), Inches(1.3), Inches(5.7), Inches(0.06), fill_color=C_WARN)
    label(s, "CPU Mode  (sequential)",
          x=Inches(0.6), y=Inches(1.38), w=Inches(5.4), h=Inches(0.45),
          size=Pt(18), bold=True, color=C_WARN)

    cpu_items = [
        "Ants construct tours one at a time",
        "Each iteration: loop over all ants",
        "Full CPU power used for one ant at a time",
        "Performance limited by clock speed",
        "50 ants processed serially per iteration",
    ]
    bullet_list(s, ["  " + x for x in cpu_items],
                x=Inches(0.6), y=Inches(1.95), w=Inches(5.4),
                size=Pt(15), color=C_TEXT, spacing=Inches(0.44))

    # GPU side
    card(s, Inches(6.6), Inches(1.3), Inches(6.2), Inches(5.6), color=C_SURFACE)
    box(s, Inches(6.6), Inches(1.3), Inches(6.2), Inches(0.06), fill_color=C_ACCENT2)
    label(s, "CUDA Mode  (parallel)",
          x=Inches(6.75), y=Inches(1.38), w=Inches(5.9), h=Inches(0.45),
          size=Pt(18), bold=True, color=C_ACCENT2)

    gpu_items = [
        "One GPU thread per ant",
        "All ants build tours simultaneously",
        "GTX 1660 Super: 1,280 CUDA cores",
        "Tour construction kernel runs in parallel",
        "Pheromone update on CPU after each iteration",
        "Memory transfers: pheromone matrix each iteration",
    ]
    bullet_list(s, ["  " + x for x in gpu_items],
                x=Inches(6.75), y=Inches(1.95), w=Inches(5.9),
                size=Pt(15), color=C_TEXT, spacing=Inches(0.44))

    # Arrow in the middle
    label(s, "vs",
          x=Inches(5.9), y=Inches(3.8), w=Inches(0.7), h=Inches(0.6),
          size=Pt(22), bold=True, color=C_SUBTEXT, align=PP_ALIGN.CENTER)


def slide_parameters(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Experiment Parameters")
    hline(s, Inches(1.1))

    params = [
        ("Cities",             "50",    "Number of cities in the problem instance"),
        ("Ants",               "50",    "Number of ants constructing tours each iteration"),
        ("Iterations",         "200",   "Number of optimization cycles"),
        ("Alpha  (alpha)",     "1.0",   "Pheromone exponent — weight given to existing trails"),
        ("Beta  (beta)",       "3.0",   "Heuristic exponent — preference for shorter edges"),
        ("Rho  (rho)",         "0.1",   "Evaporation rate — fraction of pheromone lost per iteration"),
        ("Q  (deposit)",       "100.0", "Pheromone deposit constant (proportional to tour quality)"),
        ("Init. pheromone",    "1.0",   "Starting pheromone level on all edges"),
    ]

    tbl = table_shape(s, len(params) + 1, 3,
                      x=Inches(0.55), y=Inches(1.25),
                      w=Inches(12.2), h=Inches(5.8),
                      font_size=Pt(15))

    set_cell(tbl, 0, 0, "Parameter",   bold=True)
    set_cell(tbl, 0, 1, "Value",       bold=True)
    set_cell(tbl, 0, 2, "Description", bold=True)

    for r, (name, val, desc) in enumerate(params, start=1):
        set_cell(tbl, r, 0, name, align=PP_ALIGN.LEFT)
        set_cell(tbl, r, 1, val)
        set_cell(tbl, r, 2, desc, align=PP_ALIGN.LEFT)


def slide_results_m1(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Results: Machine 1")
    hline(s, Inches(1.1))

    # Spec card
    card(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.85), color=C_SURFACE2)
    label(s, "AMD Ryzen 7 5700G  |  3801 MHz  |  8 cores / 16 threads     +     NVIDIA GeForce GTX 1660 Super  |  1,280 CUDA cores",
          x=Inches(0.7), y=Inches(1.38), w=Inches(11.9), h=Inches(0.6),
          size=Pt(15), color=C_TEXT, align=PP_ALIGN.CENTER)

    runs = ["Run 1", "Run 2", "Run 3", "Run 4", "Run 5", "Average"]
    cpu  = ["10.61 s", "10.92 s", "11.27 s", "10.21 s", "10.40 s", "10.68 s"]
    gpu  = ["1.32 s",  "1.31 s",  "1.33 s",  "1.34 s",  "1.32 s",  "1.32 s"]

    tbl = table_shape(s, 3, 7,
                      x=Inches(0.55), y=Inches(2.35),
                      w=Inches(12.2), h=Inches(1.8),
                      font_size=Pt(15))

    set_cell(tbl, 0, 0, "Mode")
    for c, r in enumerate(runs, 1):
        set_cell(tbl, 0, c, r, bold=True)

    set_cell(tbl, 1, 0, "CPU  (Ryzen 7 5700G)", align=PP_ALIGN.LEFT)
    set_cell(tbl, 2, 0, "CUDA (GTX 1660 Super)", align=PP_ALIGN.LEFT)
    for c, (cv, gv) in enumerate(zip(cpu, gpu), 1):
        bold_avg = (c == 6)
        set_cell(tbl, 1, c, cv, bold=bold_avg)
        set_cell(tbl, 2, c, gv, bold=bold_avg)

    # Speedup callout
    card(s, Inches(3.5), Inches(4.45), Inches(6.3), Inches(1.35), color=C_ACCENT)
    label(s, "Speedup:  ~8.1x  faster with CUDA",
          x=Inches(3.5), y=Inches(4.55), w=Inches(6.3), h=Inches(0.7),
          size=Pt(26), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    label(s, "10.68 s  (CPU)   ->   1.32 s  (CUDA)",
          x=Inches(3.5), y=Inches(5.2), w=Inches(6.3), h=Inches(0.45),
          size=Pt(16), color=C_BG, align=PP_ALIGN.CENTER)


def slide_results_m2(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Results: Machine 2")
    hline(s, Inches(1.1))

    card(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.85), color=C_SURFACE2)
    label(s, "Intel Core i7-13700KF  |  3400 MHz  |  16 cores / 24 threads     +     NVIDIA GeForce GTX 1660 Super  |  1,280 CUDA cores",
          x=Inches(0.7), y=Inches(1.38), w=Inches(11.9), h=Inches(0.6),
          size=Pt(15), color=C_TEXT, align=PP_ALIGN.CENTER)

    runs = ["Run 1", "Run 2", "Run 3", "Run 4", "Run 5", "Average"]
    cpu  = ["5.20 s", "5.18 s", "5.17 s", "5.22 s", "5.23 s", "5.20 s"]
    gpu  = ["0.86 s", "0.84 s", "0.86 s", "0.85 s", "0.85 s", "0.85 s"]

    tbl = table_shape(s, 3, 7,
                      x=Inches(0.55), y=Inches(2.35),
                      w=Inches(12.2), h=Inches(1.8),
                      font_size=Pt(15))

    set_cell(tbl, 0, 0, "Mode")
    for c, r in enumerate(runs, 1):
        set_cell(tbl, 0, c, r, bold=True)

    set_cell(tbl, 1, 0, "CPU  (i7-13700KF)",    align=PP_ALIGN.LEFT)
    set_cell(tbl, 2, 0, "CUDA (GTX 1660 Super)", align=PP_ALIGN.LEFT)
    for c, (cv, gv) in enumerate(zip(cpu, gpu), 1):
        bold_avg = (c == 6)
        set_cell(tbl, 1, c, cv, bold=bold_avg)
        set_cell(tbl, 2, c, gv, bold=bold_avg)

    card(s, Inches(3.5), Inches(4.45), Inches(6.3), Inches(1.35), color=C_ACCENT2)
    label(s, "Speedup:  ~6.1x  faster with CUDA",
          x=Inches(3.5), y=Inches(4.55), w=Inches(6.3), h=Inches(0.7),
          size=Pt(26), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    label(s, "5.20 s  (CPU)   ->   0.85 s  (CUDA)",
          x=Inches(3.5), y=Inches(5.2), w=Inches(6.3), h=Inches(0.45),
          size=Pt(16), color=C_BG, align=PP_ALIGN.CENTER)


def slide_comparison(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Side-by-Side Comparison")
    hline(s, Inches(1.1))

    # Summary table
    tbl = table_shape(s, 5, 4,
                      x=Inches(0.55), y=Inches(1.3),
                      w=Inches(12.2), h=Inches(2.8),
                      font_size=Pt(16))

    headers = ["", "CPU Average", "CUDA Average", "Speedup"]
    row1    = ["Machine 1  (Ryzen 7 5700G)", "10.68 s", "1.32 s", "~8.1x"]
    row2    = ["Machine 2  (i7-13700KF)",    "5.20 s",  "0.85 s", "~6.1x"]
    row3    = ["GPU used (both machines)", "—", "GTX 1660 Super  |  1,280 cores", "—"]
    row4    = ["CUDA cores vs CPU cores",  "8 / 16 cores  |  16 / 24 cores",
               "1,280 CUDA cores", "—"]

    for c, h in enumerate(headers):
        set_cell(tbl, 0, c, h, bold=True)
    for c, v in enumerate(row1):
        set_cell(tbl, 1, c, v, bold=(c == 3), align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    for c, v in enumerate(row2):
        set_cell(tbl, 2, c, v, bold=(c == 3), align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    for c, v in enumerate(row3):
        set_cell(tbl, 3, c, v, align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    for c, v in enumerate(row4):
        set_cell(tbl, 4, c, v, align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)

    label(s, "Key observations",
          x=Inches(0.55), y=Inches(4.35), w=Inches(12.2), h=Inches(0.42),
          size=Pt(17), bold=True, color=C_ACCENT)

    obs = [
        "  The same GPU (GTX 1660 Super) outperforms both CPUs by 6-8x — hardware generation matters less than parallelism.",
        "  Machine 2's faster CPU closes the absolute gap (0.85 s vs 1.32 s) but the relative speedup is similar.",
        "  CUDA time is nearly constant across runs — GPU execution is highly deterministic.",
        "  CPU time varies slightly (~1 s range on M1) due to OS scheduling and cache effects.",
    ]
    bullet_list(s, obs,
                x=Inches(0.55), y=Inches(4.85), w=Inches(12.2),
                size=Pt(15), color=C_TEXT, spacing=Inches(0.43))


def slide_conclusions(prs):
    s = blank_slide(prs)
    bg(s)
    heading(s, "Conclusions")
    hline(s, Inches(1.1))

    points = [
        ("CUDA delivers 6-8x speedup",
         "Even a mid-range GPU (GTX 1660 Super, released 2019) substantially outperforms modern multi-core CPUs for this workload."),
        ("Parallelism is the key",
         "ACO is naturally parallel: ants are independent. Mapping each ant to a GPU thread eliminates the main bottleneck."),
        ("ACO finds good solutions",
         "The algorithm consistently improves the tour over 200 iterations, with ~20-30% improvement from initial to final solution."),
        ("Scalability advantage grows",
         "With more ants or cities, the GPU advantage increases — the CPU must do proportionally more sequential work."),
    ]

    for i, (title, desc) in enumerate(points):
        col = i % 2
        row = i // 2
        cx = Inches(0.55) + col * Inches(6.2)
        cy = Inches(1.4)  + row * Inches(2.6)

        card(s, cx, cy, Inches(5.95), Inches(2.35))
        box(s, cx, cy, Inches(0.1), Inches(2.35), fill_color=C_ACCENT)
        label(s, title,
              x=cx + Inches(0.2), y=cy + Inches(0.12), w=Inches(5.6), h=Inches(0.5),
              size=Pt(17), bold=True, color=C_ACCENT)
        label(s, desc,
              x=cx + Inches(0.2), y=cy + Inches(0.68), w=Inches(5.6), h=Inches(1.5),
              size=Pt(14), color=C_TEXT)

    label(s, "CUDA parallelization makes ACO practical for large-scale TSP instances where CPU-based solutions would be prohibitively slow.",
          x=Inches(0.55), y=Inches(6.8), w=Inches(12.2), h=Inches(0.5),
          size=Pt(15), color=C_SUBTEXT, align=PP_ALIGN.CENTER)


# ── Build & save ──────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_aco_intro(prs)
    slide_aco_algorithm(prs)
    slide_cuda(prs)
    slide_parameters(prs)
    slide_results_m1(prs)
    slide_results_m2(prs)
    slide_comparison(prs)
    slide_conclusions(prs)

    out = "tsp_aco_presentation.pptx"
    prs.save(out)
    print(f"Saved  {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
